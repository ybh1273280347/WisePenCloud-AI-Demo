from pathlib import Path
from typing import List

from chat.application.document_parse.errors import (
    UnsupportedDocumentFormatError,
)
from chat.application.document_parse.models import (
    DocumentParseResult,
    ParsedPage,
    ParsedTable,
)
from chat.application.document_parse.text_utils import normalize_text
from common.logger import log_event


class OfficeNativeParser:
    def __init__(self):
        log_event("OfficeNativeParser 初始化", handler_class=type(self).__name__)

    def parse(self, path: Path, *, file_type: str) -> DocumentParseResult:
        log_event(
            "OfficeNativeParser 路由",
            path=str(path),
            file_type=file_type,
            handler_class=type(self).__name__,
        )
        if file_type == "docx":
            return self._parse_docx(path)

        if file_type == "pptx":
            return self._parse_pptx(path)

        raise UnsupportedDocumentFormatError(file_type)

    def _parse_docx(self, path: Path) -> DocumentParseResult:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc = Document(path)
        parts: List[str] = []
        tables: List[ParsedTable] = []

        for child in doc.element.body.iterchildren():
            if child.tag == qn("w:p"):
                paragraph = Paragraph(child, doc)
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
                continue

            if child.tag != qn("w:tbl"):
                continue

            table = Table(child, doc)
            rows = self._rows_from_docx_table(table)

            if rows:
                table_index = len(tables)
                parsed_table = ParsedTable(
                    table_id=f"docx_table_{table_index}",
                    source="python_docx",
                    rows=rows,
                    page_index=0,
                    metadata={},
                )
                tables.append(parsed_table)
                parts.append(f"[Table {table_index + 1}]")
                parts.extend([" | ".join(row) for row in rows])

        text = normalize_text("\n\n".join(parts))
        page = ParsedPage(
            page_index=0,
            text=text,
            page_type="document",
            tables=tables,
            metadata={"parser": "python_docx"},
        )

        result = DocumentParseResult(
            text=text,
            source=str(path),
            file_type="docx",
            pages=[page],
            tables=tables,
            metadata={"parser": "python_docx"},
            warnings=[],
        )
        log_event(
            "OfficeNativeParser parse 完成",
            path=str(path),
            file_type="docx",
            handler_class=type(self).__name__,
            backend_class=Document.__name__,
            page_count=1,
            table_count=len(tables),
            length=len(text),
        )
        return result

    def _rows_from_docx_table(self, table) -> List[List[str]]:
        rows: List[List[str]] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        return rows

    def _parse_pptx(self, path: Path) -> DocumentParseResult:
        from pptx import Presentation

        presentation = Presentation(path)
        pages: List[ParsedPage] = []
        all_tables: List[ParsedTable] = []
        text_parts: List[str] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_parts = [f"## Slide {slide_index + 1}"]
            slide_tables: List[ParsedTable] = []

            for shape in slide.shapes:
                if shape.has_table:
                    rows: List[List[str]] = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            rows.append(cells)

                    if rows:
                        table_number = len(slide_tables) + 1
                        table = ParsedTable(
                            table_id=f"pptx_slide_{slide_index}_table_{len(slide_tables)}",
                            source="python_pptx",
                            rows=rows,
                            page_index=slide_index,
                            metadata={},
                        )
                        slide_tables.append(table)
                        all_tables.append(table)
                        slide_parts.append(f"[Table {table_number}]")
                        slide_parts.extend([" | ".join(row) for row in rows])

                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)

            slide_text = normalize_text("\n".join(slide_parts))
            text_parts.append(slide_text)
            pages.append(
                ParsedPage(
                    page_index=slide_index,
                    text=slide_text,
                    page_type="slide",
                    tables=slide_tables,
                    metadata={"parser": "python_pptx"},
                )
            )

        text = normalize_text("\n\n".join(text_parts))

        result = DocumentParseResult(
            text=text,
            source=str(path),
            file_type="pptx",
            pages=pages,
            tables=all_tables,
            metadata={"parser": "python_pptx"},
            warnings=[],
        )
        log_event(
            "OfficeNativeParser parse 完成",
            path=str(path),
            file_type="pptx",
            handler_class=type(self).__name__,
            backend_class=Presentation.__name__,
            page_count=len(pages),
            table_count=len(all_tables),
            length=len(text),
        )
        return result
